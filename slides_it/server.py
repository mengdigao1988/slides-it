from __future__ import annotations

import asyncio
import base64
import json
import logging
import mimetypes
import os
import re
import shutil
import threading
import pathlib
import signal
import subprocess
import tempfile
import urllib.error
import urllib.request
from contextlib import asynccontextmanager
from typing import List

import uvicorn
from fastapi import FastAPI, File, HTTPException, Query, Response, UploadFile
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import StreamingResponse
from fastapi.staticfiles import StaticFiles
from pydantic import BaseModel

from slides_it import __version__
from slides_it.designs import DesignManager
from slides_it.industries import IndustryManager

# ---------------------------------------------------------------------------
# Constants
# ---------------------------------------------------------------------------

OPENCODE_PORT = 4096

# Known providers that ship native support in opencode (no custom npm package needed)
_NATIVE_PROVIDERS = {"anthropic", "openai", "opencode", "openrouter", "deepseek", ""}

# Directories that are never shown in the workspace file tree.
_LS_IGNORE: set[str] = {
    ".git", ".vscode", ".venv", ".idea",
    "node_modules", "__pycache__", ".DS_Store",
    ".mypy_cache", ".pytest_cache", ".ruff_cache",
    ".tox", "dist", ".next", ".nuxt",
    ".slides-it",   # internal slides-it state — not useful to the user
}

# MCP config block for open-webSearch — always injected into opencode.json
_MCP_WEB_SEARCH_BLOCK: dict = {
    "type": "local",
    "command": ["open-websearch"],
    "environment": {
        "MODE": "stdio",
        "DEFAULT_SEARCH_ENGINE": "bing",
        "ALLOWED_SEARCH_ENGINES": "bing,baidu,brave,duckduckgo,startpage",
    },
}

# ---------------------------------------------------------------------------
# State
# ---------------------------------------------------------------------------

_opencode_proc: subprocess.Popen[bytes] | None = None
_workspace_dir: str = ""


# ---------------------------------------------------------------------------
# Lifespan
# ---------------------------------------------------------------------------

@asynccontextmanager
async def lifespan(app: FastAPI):  # type: ignore[type-arg]
    # On startup: clean up any stale slides-it AGENTS.md block left by an older
    # version that used write_rules(). slides-it no longer owns that file.
    _cleanup_stale_agents_md()
    yield
    # Cleanup on shutdown
    _stop_opencode()


# ---------------------------------------------------------------------------
# App
# ---------------------------------------------------------------------------

app = FastAPI(title="slides-it", lifespan=lifespan)

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_methods=["*"],
    allow_headers=["*"],
)


# ---------------------------------------------------------------------------
# Models
# ---------------------------------------------------------------------------

class DirEntry(BaseModel):
    name: str
    path: str
    has_children: bool


class FsEntry(BaseModel):
    name: str
    path: str
    type: str          # "directory" | "file"
    has_children: bool  # directories only; always False for files


class StartRequest(BaseModel):
    directory: str


class StatusResponse(BaseModel):
    ready: bool
    workspace: str
    version: str


class SettingsResponse(BaseModel):
    providerID: str
    apiKeyMasked: str   # "sk-ant-...abcd" or ""
    baseURL: str
    customModel: str


class SettingsRequest(BaseModel):
    providerID: str
    apiKey: str         # empty string = keep existing
    baseURL: str
    customModel: str


class DesignEntry(BaseModel):
    name: str
    description: str
    author: str
    version: str
    active: bool
    has_preview: bool


class DesignDetail(BaseModel):
    name: str
    description: str
    author: str
    version: str
    active: bool
    has_preview: bool
    skill_md: str
    preview_html: str | None


class SessionRequest(BaseModel):
    session_id: str
    parent_session_id: str = ""     # set during replay to link child → parent
    messages: list[dict] = []       # serialised ChatMessage[] from the frontend


class InstallDesignRequest(BaseModel):
    # Mode A: install from external source (URL, github:user/repo, registry name, local path)
    source: str = ""

    # Mode B: install from inline content (used by the AI agent via curl)
    # All three fields are required together when source is empty.
    name: str = ""           # kebab-case design name, e.g. "blue-minimal"
    description: str = ""
    skill_md: str = ""       # full SKILL.md content
    preview_html: str = ""   # full preview.html content (may be empty)
    activate: bool = False   # if True, activate this design after installing


class IndustryEntry(BaseModel):
    name: str
    description: str
    author: str
    version: str
    active: bool


class IndustryDetail(BaseModel):
    name: str
    description: str
    author: str
    version: str
    active: bool
    skill_md: str


class InstallIndustryRequest(BaseModel):
    # Mode A: install from external source (URL, github:user/repo, local path)
    source: str = ""

    # Mode B: install from inline content
    name: str = ""           # kebab-case industry name, e.g. "deeptech-investment"
    description: str = ""
    skill_md: str = ""       # full INDUSTRY.md body content
    activate: bool = False


class DocumentEntry(BaseModel):
    name: str          # filename only
    path: str          # relative path from workspace root
    type: str          # "pdf" | "xlsx" | "xls" | "docx" | "doc" | "pptx" | "ppt" | "csv"
    size: int          # file size in bytes
    size_human: str    # e.g. "2.3 MB"


class ExtractRequest(BaseModel):
    path: str                    # relative or absolute path to the file
    max_chars: int = 30_000      # max characters to return (server caps at 50k)
    pages: str = ""              # page range for PDF/PPTX, e.g. "1-10" (empty = all)


class ExtractResponse(BaseModel):
    path: str
    type: str
    content: str                 # extracted markdown text
    total_pages: int | None = None       # PDF/PPTX page count
    extracted_pages: str | None = None   # e.g. "1-10"
    total_sheets: int | None = None      # Excel sheet count
    sheet_names: list[str] | None = None # Excel sheet names
    truncated: bool = False
    hint: str = ""               # guidance for AI on how to get more content


class DocumentInfoResponse(BaseModel):
    path: str
    type: str
    size: int
    size_human: str
    total_pages: int | None = None
    total_sheets: int | None = None
    sheet_names: list[str] | None = None


class FileRenameRequest(BaseModel):
    path: str       # absolute path of the file/directory to rename
    new_name: str   # new filename only (not a full path), no path separators


class FileCreateRequest(BaseModel):
    name: str       # filename only (no path separators), created in workspace root


class DirCreateRequest(BaseModel):
    name: str       # directory name only (no path separators), created in workspace root


class FileContentRequest(BaseModel):
    path: str       # absolute path of the file to write
    content: str    # full file content (text)


class BundleRequest(BaseModel):
    path: str       # absolute path of the HTML file to bundle


class ReplayRequest(BaseModel):
    session_id: str
    provider_id: str = ""   # empty = auto-detect from active model
    model_id: str = ""      # empty = auto-detect from active model


# ---------------------------------------------------------------------------
# API routes
# ---------------------------------------------------------------------------

@app.get("/api/dirs", response_model=list[DirEntry])
def list_dirs(path: str = Query(default="~")) -> list[DirEntry]:
    """
    List immediate subdirectories of `path`.
    Returns only directories, sorted: hidden last, then alpha.
    """
    resolved = pathlib.Path(path).expanduser().resolve()
    if not resolved.exists() or not resolved.is_dir():
        raise HTTPException(status_code=404, detail=f"Directory not found: {path}")

    entries: list[DirEntry] = []
    try:
        for child in sorted(resolved.iterdir(), key=lambda p: (p.name.startswith("."), p.name.lower())):
            if not child.is_dir():
                continue
            # Check if it has any subdirectory children (for chevron indicator)
            try:
                has_children = any(c.is_dir() for c in child.iterdir())
            except PermissionError:
                has_children = False
            entries.append(DirEntry(name=child.name, path=str(child), has_children=has_children))
    except PermissionError:
        raise HTTPException(status_code=403, detail=f"Permission denied: {path}")

    return entries


@app.get("/api/ls", response_model=list[FsEntry])
def list_entries(path: str = Query(default="~")) -> list[FsEntry]:
    """
    List all immediate children (directories AND files) of `path`.
    Directories come first, then files, both sorted alphabetically.
    Hidden entries (dot-prefixed) come last within each group.
    """
    resolved = pathlib.Path(path).expanduser().resolve()
    if not resolved.exists() or not resolved.is_dir():
        raise HTTPException(status_code=404, detail=f"Directory not found: {path}")

    entries: list[FsEntry] = []
    try:
        children = sorted(
            resolved.iterdir(),
            key=lambda p: (p.is_file(), p.name.startswith("."), p.name.lower()),
        )
        for child in children:
            # Skip noise directories / files that users don't need to see
            if child.name in _LS_IGNORE:
                continue
            if child.is_symlink():
                # Resolve symlinks to get real type
                try:
                    real = child.resolve()
                    is_dir = real.is_dir()
                except OSError:
                    continue
            else:
                is_dir = child.is_dir()

            has_children = False
            if is_dir:
                try:
                    has_children = any(True for _ in child.iterdir())
                except PermissionError:
                    pass

            entries.append(FsEntry(
                name=child.name,
                path=str(child),
                type="directory" if is_dir else "file",
                has_children=has_children,
            ))
    except PermissionError:
        raise HTTPException(status_code=403, detail=f"Permission denied: {path}")

    return entries


_IGNORE_CONTENT = """\
# slides-it: exclude binary/media files from AI file search (ripgrep respects this file
# regardless of whether the directory is a git repository).

# Images
*.png
*.jpg
*.jpeg
*.gif
*.webp
*.bmp
*.ico
*.tiff
*.tif
*.avif

# Video
*.mp4
*.mov
*.avi
*.mkv
*.webm
*.m4v
*.wmv

# Audio
*.mp3
*.wav
*.ogg
*.flac
*.aac
*.m4a

# Documents
*.pdf
*.docx
*.xlsx
*.pptx
*.doc
*.xls
*.ppt

# Archives
*.zip
*.tar
*.gz
*.bz2
*.7z
*.rar
*.tgz

# Fonts
*.woff
*.woff2
*.ttf
*.otf
*.eot

# Database / compiled binary
*.db
*.sqlite
*.sqlite3
*.bin
*.exe
*.dll
*.so
*.dylib

# Hidden directories (dot-prefixed)
.*/
"""


def _ensure_ignore_file(directory: pathlib.Path) -> None:
    """
    Write a .ignore file in the workspace root if one does not already exist.

    ripgrep (used by OpenCode's grep/glob/list tools) reads .ignore unconditionally,
    unlike .gitignore which requires a git repository. This prevents the AI from
    accidentally scanning binary or media files.
    """
    ignore_path = directory / ".ignore"
    if not ignore_path.exists():
        ignore_path.write_text(_IGNORE_CONTENT, encoding="utf-8")


@app.post("/api/start")
def start_workspace(req: StartRequest) -> dict[str, str]:
    """
    Start opencode serve in the given directory.
    If opencode is already healthy on the expected port, reuse it.
    """
    global _opencode_proc, _workspace_dir

    directory = pathlib.Path(req.directory).expanduser().resolve()
    if not directory.exists() or not directory.is_dir():
        raise HTTPException(status_code=400, detail=f"Directory not found: {req.directory}")

    # Create .slides-it directory in workspace (for future session persistence etc.)
    (directory / ".slides-it").mkdir(exist_ok=True)

    # Write .ignore so ripgrep skips binary/media files even without a git repo
    _ensure_ignore_file(directory)

    # Ensure MCP web-search config is present in opencode.json
    _ensure_mcp_config(str(directory))

    # If opencode is already healthy, just update workspace and reuse it
    if _is_opencode_healthy():
        _workspace_dir = str(directory)
        return {"status": "starting", "workspace": str(directory)}

    # Stop any existing managed opencode process, then free the port
    _stop_opencode()

    # Start opencode serve
    _opencode_proc = subprocess.Popen(
        ["opencode", "serve", "--port", str(OPENCODE_PORT)],
        cwd=str(directory),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
    )
    _workspace_dir = str(directory)

    return {"status": "starting", "workspace": str(directory)}


@app.get("/api/status", response_model=StatusResponse)
def get_status() -> StatusResponse:
    """Check whether opencode serve is up and healthy."""
    # First check our managed process is alive (if we started one)
    if _opencode_proc and _opencode_proc.poll() is not None:
        return StatusResponse(ready=False, workspace=_workspace_dir, version=__version__)

    # Then actually ping opencode regardless of how it was started
    try:
        with urllib.request.urlopen(
            f"http://localhost:{OPENCODE_PORT}/global/health", timeout=2
        ) as resp:
            data = json.loads(resp.read())
            return StatusResponse(
                ready=data.get("healthy", False),
                workspace=_workspace_dir,
                version=__version__,
            )
    except Exception:
        return StatusResponse(ready=False, workspace=_workspace_dir, version=__version__)


@app.post("/api/shutdown")
async def shutdown() -> dict[str, str]:
    """
    Graceful shutdown triggered by the browser window closing.
    Sends SIGTERM to self after the response is delivered.
    """
    asyncio.get_running_loop().call_later(0.1, lambda: os.kill(os.getpid(), signal.SIGTERM))
    return {"status": "shutting_down"}


@app.get("/api/session")
def get_session() -> dict[str, object]:
    """
    Return the persisted session state for the current workspace.

    Walks the parent chain to reconstruct the full message history:

        session-C.json  { parent: "ses_B", messages: [C's msgs] }
        session-B.json  { parent: "ses_A", messages: [B's msgs] }
        session-A.json  { messages: [A's msgs] }

    Returns:
        session_id:       current session ID (from the ``current`` pointer)
        messages:         full reconstructed message history (A + compact + B + compact + C)
        recent_messages:  only the current session's own messages (C's msgs)
                          — used by the frontend to inject context on restart
    """
    if not _workspace_dir:
        return {"session_id": None, "messages": [], "recent_messages": []}
    slides_dir = pathlib.Path(_workspace_dir) / ".slides-it"
    slides_dir.mkdir(exist_ok=True)
    current_file = slides_dir / "current"
    if not current_file.exists():
        return {"session_id": None, "messages": [], "recent_messages": []}
    try:
        session_id = current_file.read_text(encoding="utf-8").strip()
        if not session_id:
            return {"session_id": None, "messages": [], "recent_messages": []}

        # Walk the parent chain and collect each session's messages
        chain: list[tuple[str, list[dict]]] = []  # [(sid, msgs), ...]
        visited: set[str] = set()
        sid: str | None = session_id
        while sid and sid not in visited:
            visited.add(sid)
            sf = slides_dir / f"session-{sid}.json"
            if not sf.exists():
                break
            data = json.loads(sf.read_text(encoding="utf-8"))
            chain.append((sid, data.get("messages", [])))
            sid = data.get("parent") or None

        # chain is [current, parent, grandparent, ...] — reverse to get chronological order
        chain.reverse()

        # Build full message history with compact separators between sessions
        all_messages: list[dict] = []
        for i, (_sid, msgs) in enumerate(chain):
            if i > 0 and msgs:
                # Insert a compact separator between parent and child sessions
                all_messages.append({
                    "id": f"compact-{_sid}",
                    "role": "system",
                    "text": "--- Context compacted — conversation continues ---",
                    "streaming": False,
                    "error": None,
                    "timestamp": msgs[0].get("timestamp") if msgs else None,
                    "tools": [],
                    "compact": True,
                })
            all_messages.extend(msgs)

        # recent_messages = current session's own messages only
        recent_messages = chain[-1][1] if chain else []

        return {
            "session_id": session_id,
            "messages": all_messages,
            "recent_messages": recent_messages,
        }
    except Exception:
        return {"session_id": None, "messages": [], "recent_messages": []}


@app.put("/api/session")
def save_session(req: SessionRequest) -> dict[str, str]:
    """
    Persist messages to <workspace>/.slides-it/session-<id>.json
    and update the 'current' pointer.

    Each session file stores only its own messages.  The optional
    ``parent_session_id`` field links to the previous session (set during
    replay).  Old session files are kept (not deleted).
    """
    if not _workspace_dir:
        raise HTTPException(status_code=400, detail="No active workspace")
    slides_dir = pathlib.Path(_workspace_dir) / ".slides-it"
    slides_dir.mkdir(exist_ok=True)

    # Build the session data — only this session's messages
    session_data: dict[str, object] = {"messages": req.messages}
    if req.parent_session_id:
        session_data["parent"] = req.parent_session_id

    session_file = slides_dir / f"session-{req.session_id}.json"

    # If the file already exists and has a parent, preserve the parent link
    # even if this save call doesn't include it (e.g. session.idle auto-save)
    if not req.parent_session_id and session_file.exists():
        try:
            existing = json.loads(session_file.read_text(encoding="utf-8"))
            if existing.get("parent"):
                session_data["parent"] = existing["parent"]
        except Exception:
            pass

    session_file.write_text(
        json.dumps(session_data, ensure_ascii=False),
        encoding="utf-8",
    )
    # Update the current pointer
    (slides_dir / "current").write_text(req.session_id, encoding="utf-8")
    return {"status": "saved"}


@app.get("/api/models")
def get_models() -> dict[str, object]:
    """
    Return the list of available models (via `opencode models` CLI)
    and the currently active model ID.
    """
    dm = DesignManager()
    current = dm.get_model()

    try:
        result = subprocess.run(
            ["opencode", "models"],
            capture_output=True, text=True, timeout=10,
            cwd=_workspace_dir if _workspace_dir else None,
        )
        models = [line.strip() for line in result.stdout.splitlines() if line.strip()]
    except Exception:
        models = []

    return {"models": models, "current": current}


@app.put("/api/model")
def set_model(body: dict[str, str]) -> dict[str, str]:
    """Save the active model ID to slides-it config."""
    model_id = body.get("modelID", "").strip()
    if not model_id:
        raise HTTPException(status_code=400, detail="modelID is required")
    DesignManager().set_model(model_id)
    return {"modelID": model_id}


@app.get("/api/design/{name}/skill")
def get_design_skill(name: str, industry: str = Query(default="")) -> dict[str, str]:
    """
    Return the combined system prompt for the given design.

    Concatenates core SKILL.md + industry INDUSTRY.md + design DESIGN.md.
    The frontend sends this as the `system` field in POST /session/:id/prompt_async
    so that the active design's visual style and industry context are injected
    on every message without touching any config files on disk.

    Args:
        name: Design name.
        industry: Industry name (optional, defaults to active industry).
    """
    dm = DesignManager()
    industry_name = industry if industry else None
    try:
        skill = dm.build_prompt(name, industry_name=industry_name)
    except ValueError as e:
        raise HTTPException(status_code=404, detail=str(e))
    return {"skill": skill}


@app.get("/api/design/{name}/preview")
def get_design_preview(name: str) -> dict[str, str]:
    """
    Return the raw HTML content of preview.html for the given design.
    The frontend injects this into an iframe via srcdoc.
    """
    dm = DesignManager()
    path = dm._design_path(name)
    if not path:
        raise HTTPException(status_code=404, detail=f"Design '{name}' not found")
    preview_file = path / "preview.html"
    if not preview_file.exists():
        raise HTTPException(status_code=404, detail=f"Design '{name}' has no preview.html")
    return {"html": preview_file.read_text(encoding="utf-8")}


@app.get("/api/design/{name}", response_model=DesignDetail)
def get_design(name: str) -> DesignDetail:
    """
    Return full details for a single design — metadata, skill text, and preview.html.

    Used by the AI agent to fetch the active design's visual reference in one
    call before generating slides. Also available to the frontend for any future
    use that needs all fields together.

    preview_html is null if the design has no preview.html.
    skill_md contains the skill text body from DESIGN.md (not the combined prompt).
    """
    dm = DesignManager()
    path = dm._design_path(name)
    if not path:
        raise HTTPException(status_code=404, detail=f"Design '{name}' not found")
    info = dm._parse_design_file(path / "DESIGN.md")
    if not info:
        raise HTTPException(status_code=404, detail=f"Design '{name}' has no DESIGN.md")
    preview_file = path / "preview.html"
    return DesignDetail(
        name=info.name,
        description=info.description,
        author=info.author,
        version=info.version,
        active=info.name == dm.active(),
        has_preview=preview_file.exists(),
        skill_md=info.skill_text,
        preview_html=preview_file.read_text(encoding="utf-8") if preview_file.exists() else None,
    )


@app.get("/api/designs", response_model=list[DesignEntry])
def list_designs() -> list[DesignEntry]:
    """Return all installed designs with metadata."""
    dm = DesignManager()
    active = dm.active()
    result = []
    for info in dm.list():
        path = dm._design_path(info.name)
        has_preview = bool(path and (path / "preview.html").exists())
        result.append(DesignEntry(
            name=info.name,
            description=info.description,
            author=info.author,
            version=info.version,
            active=info.name == active,
            has_preview=has_preview,
        ))
    return result


@app.post("/api/designs/install")
def install_design(req: InstallDesignRequest) -> dict[str, str]:
    """
    Install a design — two modes:

    Mode A (source): install from any external source.
        { "source": "https://...", "activate": true }
        { "source": "github:user/repo" }
        { "source": "dark-neon" }          ← registry lookup

    Mode B (inline): install from content generated by the AI agent.
        {
          "name": "blue-minimal",
          "description": "Clean blue theme with minimal decoration",
          "skill_md": "## Visual Style...",
          "preview_html": "<!DOCTYPE html>...",
          "activate": true
        }

    In Mode B the server writes a temporary directory containing DESIGN.md
    and optionally preview.html, then installs it via the same
    _install_from_path() codepath used by Mode A.  This keeps all install
    logic in one place.
    """
    dm = DesignManager()

    if req.source.strip():
        # ── Mode A: external source ───────────────────────────────────────
        source = req.source.strip()
        try:
            name = dm.install(source)
        except Exception as e:
            raise HTTPException(status_code=400, detail=str(e))
        if req.activate:
            try:
                dm.activate(name)
            except ValueError as e:
                raise HTTPException(status_code=400, detail=str(e))
        return {"name": name, "status": "installed", "activated": str(req.activate).lower()}

    else:
        # ── Mode B: inline content from agent ────────────────────────────
        name = req.name.strip()
        if not name:
            raise HTTPException(status_code=400, detail="'name' is required when 'source' is empty")
        if not req.skill_md.strip():
            raise HTTPException(status_code=400, detail="'skill_md' is required when 'source' is empty")

        # Validate name: kebab-case, no path traversal
        import re as _re
        if not _re.match(r'^[a-z0-9]+(?:-[a-z0-9]+)*$', name):
            raise HTTPException(
                status_code=400,
                detail="'name' must be kebab-case (lowercase letters, digits, hyphens only)",
            )

        with tempfile.TemporaryDirectory() as tmp:
            tmp_path = pathlib.Path(tmp)

            # Write DESIGN.md — YAML frontmatter + skill text body in one file
            design_md = (
                "---\n"
                f"name: {name}\n"
                f"description: {req.description.strip() or 'AI-generated design'}\n"
                "author: ai-generated\n"
                "version: 1.0.0\n"
                "---\n\n"
                f"{req.skill_md}"
            )
            (tmp_path / "DESIGN.md").write_text(design_md, encoding="utf-8")

            # Write preview.html if provided
            if req.preview_html.strip():
                (tmp_path / "preview.html").write_text(req.preview_html, encoding="utf-8")

            try:
                installed_name = dm.install(str(tmp_path))
            except Exception as e:
                raise HTTPException(status_code=400, detail=str(e))

        if req.activate:
            try:
                dm.activate(installed_name)
            except ValueError as e:
                raise HTTPException(status_code=400, detail=str(e))

        return {"name": installed_name, "status": "installed", "activated": str(req.activate).lower()}


@app.delete("/api/designs/{name}")
def delete_design(name: str) -> dict[str, str]:
    """Remove an installed design."""
    dm = DesignManager()
    try:
        dm.remove(name)
    except ValueError as e:
        raise HTTPException(status_code=400, detail=str(e))
    return {"name": name, "status": "removed"}


@app.put("/api/designs/{name}/activate")
def activate_design(name: str) -> dict[str, str]:
    """Set the active design."""
    dm = DesignManager()
    try:
        dm.activate(name)
    except ValueError as e:
        raise HTTPException(status_code=404, detail=str(e))
    return {"name": name, "status": "activated"}


# ---------------------------------------------------------------------------
# Industry API routes
# ---------------------------------------------------------------------------

@app.get("/api/industries", response_model=list[IndustryEntry])
def list_industries() -> list[IndustryEntry]:
    """Return all installed industries with metadata."""
    im = IndustryManager()
    active = im.active()
    result = []
    for info in im.list():
        result.append(IndustryEntry(
            name=info.name,
            description=info.description,
            author=info.author,
            version=info.version,
            active=info.name == active,
        ))
    return result


@app.get("/api/industry/{name}", response_model=IndustryDetail)
def get_industry(name: str) -> IndustryDetail:
    """Return full details for a single industry — metadata and skill text."""
    im = IndustryManager()
    path = im._industry_path(name)
    if not path:
        raise HTTPException(status_code=404, detail=f"Industry '{name}' not found")
    info = im._parse_industry_file(path / "INDUSTRY.md")
    if not info:
        raise HTTPException(status_code=404, detail=f"Industry '{name}' has no INDUSTRY.md")
    return IndustryDetail(
        name=info.name,
        description=info.description,
        author=info.author,
        version=info.version,
        active=info.name == im.active(),
        skill_md=info.skill_text,
    )


@app.get("/api/industry/{name}/skill")
def get_industry_skill(name: str) -> dict[str, str]:
    """Return the raw skill text body from INDUSTRY.md for a given industry."""
    im = IndustryManager()
    try:
        skill = im.get_skill_md(name)
    except ValueError as e:
        raise HTTPException(status_code=404, detail=str(e))
    return {"skill": skill}


@app.put("/api/industries/{name}/activate")
def activate_industry(name: str) -> dict[str, str]:
    """Set the active industry."""
    im = IndustryManager()
    try:
        im.activate(name)
    except ValueError as e:
        raise HTTPException(status_code=404, detail=str(e))
    return {"name": name, "status": "activated"}


@app.post("/api/industries/install")
def install_industry(req: InstallIndustryRequest) -> dict[str, str]:
    """
    Install an industry — two modes:

    Mode A (source): install from any external source.
        { "source": "https://...", "activate": true }
        { "source": "github:user/repo" }

    Mode B (inline): install from content generated by the AI agent or user.
        {
          "name": "biotech-pharma",
          "description": "Biotech & pharmaceutical industry",
          "skill_md": "## Report Structure...",
          "activate": true
        }
    """
    im = IndustryManager()

    if req.source.strip():
        # ── Mode A: external source ───────────────────────────────────────
        source = req.source.strip()
        try:
            name = im.install(source)
        except Exception as e:
            raise HTTPException(status_code=400, detail=str(e))
        if req.activate:
            try:
                im.activate(name)
            except ValueError as e:
                raise HTTPException(status_code=400, detail=str(e))
        return {"name": name, "status": "installed", "activated": str(req.activate).lower()}

    else:
        # ── Mode B: inline content ───────────────────────────────────────
        name = req.name.strip()
        if not name:
            raise HTTPException(status_code=400, detail="'name' is required when 'source' is empty")
        if not req.skill_md.strip():
            raise HTTPException(status_code=400, detail="'skill_md' is required when 'source' is empty")

        # Validate name: kebab-case, no path traversal
        import re as _re
        if not _re.match(r'^[a-z0-9]+(?:-[a-z0-9]+)*$', name):
            raise HTTPException(
                status_code=400,
                detail="'name' must be kebab-case (lowercase letters, digits, hyphens only)",
            )

        with tempfile.TemporaryDirectory() as tmp:
            tmp_path = pathlib.Path(tmp)

            # Write INDUSTRY.md — YAML frontmatter + skill text body in one file
            industry_md = (
                "---\n"
                f"name: {name}\n"
                f"description: {req.description.strip() or 'User-defined industry'}\n"
                "author: user\n"
                "version: 1.0.0\n"
                "---\n\n"
                f"{req.skill_md}"
            )
            (tmp_path / "INDUSTRY.md").write_text(industry_md, encoding="utf-8")

            try:
                installed_name = im.install(str(tmp_path))
            except Exception as e:
                raise HTTPException(status_code=400, detail=str(e))

        if req.activate:
            try:
                im.activate(installed_name)
            except ValueError as e:
                raise HTTPException(status_code=400, detail=str(e))

        return {"name": installed_name, "status": "installed", "activated": str(req.activate).lower()}


@app.delete("/api/industries/{name}")
def delete_industry(name: str) -> dict[str, str]:
    """Remove an installed industry."""
    im = IndustryManager()
    try:
        im.remove(name)
    except ValueError as e:
        raise HTTPException(status_code=400, detail=str(e))
    return {"name": name, "status": "removed"}


# ---------------------------------------------------------------------------
# Document extraction helpers
# ---------------------------------------------------------------------------

# File extensions the document API can discover and extract.
_DOCUMENT_EXTENSIONS: set[str] = {
    ".pdf", ".xlsx", ".xls", ".docx", ".doc", ".pptx", ".ppt", ".csv",
}

# Image extensions — discoverable but not extractable in the same way.
_IMAGE_EXTENSIONS: set[str] = {
    ".png", ".jpg", ".jpeg", ".gif", ".webp", ".bmp", ".tiff", ".tif", ".avif",
}

# Absolute ceiling — even if the caller asks for more, never exceed this.
_MAX_EXTRACT_CHARS = 50_000


def _human_size(size_bytes: int) -> str:
    """Return a human-readable file size string."""
    if size_bytes < 1024:
        return f"{size_bytes} B"
    elif size_bytes < 1024 * 1024:
        return f"{size_bytes / 1024:.1f} KB"
    elif size_bytes < 1024 * 1024 * 1024:
        return f"{size_bytes / (1024 * 1024):.1f} MB"
    return f"{size_bytes / (1024 * 1024 * 1024):.1f} GB"


def _parse_page_range(pages_str: str, total: int) -> tuple[int, int]:
    """Parse a page range string like '1-10' into (start, end) 0-indexed inclusive.

    Returns (0, total-1) if the string is empty or invalid.
    """
    if not pages_str.strip():
        return 0, total - 1
    try:
        parts = pages_str.strip().split("-")
        if len(parts) == 1:
            p = max(1, int(parts[0]))
            return p - 1, p - 1
        start = max(1, int(parts[0]))
        end = min(total, int(parts[1]))
        return start - 1, end - 1
    except (ValueError, IndexError):
        return 0, total - 1


def _extract_pdf(file_path: pathlib.Path, max_chars: int, pages: str) -> ExtractResponse:
    """Extract text and tables from a PDF file using pdfplumber."""
    import pdfplumber

    with pdfplumber.open(file_path) as pdf:
        total_pages = len(pdf.pages)
        start, end = _parse_page_range(pages, total_pages)
        end = min(end, total_pages - 1)

        parts: list[str] = []
        char_count = 0
        actual_end = start

        for i in range(start, end + 1):
            page = pdf.pages[i]
            page_text = page.extract_text() or ""

            # Also extract tables as markdown
            tables = page.extract_tables()
            table_md = ""
            for table in tables:
                if not table:
                    continue
                # Build markdown table
                rows: list[str] = []
                for ri, row in enumerate(table):
                    cells = [str(c).strip() if c else "" for c in row]
                    rows.append("| " + " | ".join(cells) + " |")
                    if ri == 0:
                        rows.append("| " + " | ".join("---" for _ in cells) + " |")
                table_md += "\n".join(rows) + "\n\n"

            section = f"## Page {i + 1}\n\n{page_text}"
            if table_md:
                section += f"\n\n### Tables\n\n{table_md}"

            if char_count + len(section) > max_chars:
                # Include as much of this page as fits
                remaining = max_chars - char_count
                if remaining > 100:
                    parts.append(section[:remaining])
                actual_end = i
                break

            parts.append(section)
            char_count += len(section)
            actual_end = i

        content = "\n\n".join(parts)
        truncated = actual_end < total_pages - 1 or char_count >= max_chars
        extracted_range = f"{start + 1}-{actual_end + 1}"
        hint = ""
        if truncated and actual_end < total_pages - 1:
            hint = f"Content truncated. Use pages=\"{actual_end + 2}-{total_pages}\" to read the rest ({total_pages - actual_end - 1} pages remaining)."

        return ExtractResponse(
            path=str(file_path),
            type="pdf",
            content=content,
            total_pages=total_pages,
            extracted_pages=extracted_range,
            truncated=truncated,
            hint=hint,
        )


def _extract_xlsx(file_path: pathlib.Path, max_chars: int) -> ExtractResponse:
    """Extract data from an Excel file as markdown tables using openpyxl."""
    import openpyxl

    wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
    sheet_names = wb.sheetnames
    parts: list[str] = []
    char_count = 0
    truncated = False

    for sheet_name in sheet_names:
        ws = wb[sheet_name]
        rows_data: list[list[str]] = []
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            rows_data.append(cells)

        if not rows_data:
            parts.append(f"## Sheet: {sheet_name}\n\n*(empty)*")
            continue

        # Build markdown table
        md_rows: list[str] = []
        for ri, cells in enumerate(rows_data):
            md_rows.append("| " + " | ".join(cells) + " |")
            if ri == 0:
                md_rows.append("| " + " | ".join("---" for _ in cells) + " |")

        section = f"## Sheet: {sheet_name}\n\n" + "\n".join(md_rows) + "\n"

        if char_count + len(section) > max_chars:
            # Include partial sheet info
            remaining = max_chars - char_count
            if remaining > 200:
                # Show header + first N rows that fit
                parts.append(section[:remaining])
            truncated = True
            break

        parts.append(section)
        char_count += len(section)

    wb.close()

    content = "\n\n".join(parts)
    hint = ""
    if truncated:
        hint = f"Content truncated at {max_chars} characters. File has {len(sheet_names)} sheet(s): {', '.join(sheet_names)}."

    return ExtractResponse(
        path=str(file_path),
        type="xlsx",
        content=content,
        total_sheets=len(sheet_names),
        sheet_names=sheet_names,
        truncated=truncated,
        hint=hint,
    )


def _extract_docx(file_path: pathlib.Path, max_chars: int) -> ExtractResponse:
    """Extract text and tables from a Word document using python-docx."""
    import docx

    doc = docx.Document(file_path)
    parts: list[str] = []
    char_count = 0
    truncated = False

    for element in doc.element.body:
        if char_count >= max_chars:
            truncated = True
            break

        tag = element.tag.split("}")[-1]  # strip namespace

        if tag == "p":
            # Paragraph
            from docx.oxml.ns import qn
            style_el = element.find(qn("w:pPr"))
            text = element.text or ""
            # Collect all run texts
            full_text = "".join(
                node.text or "" for node in element.iter()
                if node.tag.endswith("}t")
            )
            if not full_text.strip():
                continue

            # Detect heading level
            heading_level = 0
            if style_el is not None:
                pstyle = style_el.find(qn("w:pStyle"))
                if pstyle is not None:
                    val = pstyle.get(qn("w:val"), "")
                    if val.startswith("Heading"):
                        try:
                            heading_level = int(val.replace("Heading", "").strip())
                        except ValueError:
                            heading_level = 1

            if heading_level > 0:
                line = "#" * heading_level + " " + full_text
            else:
                line = full_text

            if char_count + len(line) > max_chars:
                remaining = max_chars - char_count
                if remaining > 50:
                    parts.append(line[:remaining])
                truncated = True
                break

            parts.append(line)
            char_count += len(line)

        elif tag == "tbl":
            # Table — use python-docx Table object
            from docx.table import Table as DocxTable
            tbl = DocxTable(element, doc)
            md_rows: list[str] = []
            for ri, row in enumerate(tbl.rows):
                cells = [cell.text.strip() for cell in row.cells]
                md_rows.append("| " + " | ".join(cells) + " |")
                if ri == 0:
                    md_rows.append("| " + " | ".join("---" for _ in cells) + " |")
            table_md = "\n".join(md_rows)

            if char_count + len(table_md) > max_chars:
                truncated = True
                break

            parts.append(table_md)
            char_count += len(table_md)

    content = "\n\n".join(parts)
    hint = ""
    if truncated:
        hint = f"Content truncated at {max_chars} characters."

    return ExtractResponse(
        path=str(file_path),
        type="docx",
        content=content,
        truncated=truncated,
        hint=hint,
    )


def _extract_pptx(file_path: pathlib.Path, max_chars: int, pages: str) -> ExtractResponse:
    """Extract text from a PowerPoint file using python-pptx."""
    from pptx import Presentation

    prs = Presentation(file_path)
    total_slides = len(prs.slides)
    start, end = _parse_page_range(pages, total_slides)
    end = min(end, total_slides - 1)

    parts: list[str] = []
    char_count = 0
    actual_end = start

    for i, slide in enumerate(prs.slides):
        if i < start:
            continue
        if i > end:
            break

        slide_texts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        slide_texts.append(text)

            # Extract tables
            if shape.has_table:
                tbl = shape.table
                md_rows: list[str] = []
                for ri, row in enumerate(tbl.rows):
                    cells = [cell.text.strip() for cell in row.cells]
                    md_rows.append("| " + " | ".join(cells) + " |")
                    if ri == 0:
                        md_rows.append("| " + " | ".join("---" for _ in cells) + " |")
                slide_texts.append("\n".join(md_rows))

        # Include slide notes if present
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                slide_texts.append(f"\n*Notes: {notes}*")

        section = f"## Slide {i + 1}\n\n" + "\n\n".join(slide_texts)

        if char_count + len(section) > max_chars:
            remaining = max_chars - char_count
            if remaining > 100:
                parts.append(section[:remaining])
            actual_end = i
            break

        parts.append(section)
        char_count += len(section)
        actual_end = i

    content = "\n\n".join(parts)
    truncated = actual_end < total_slides - 1 or char_count >= max_chars
    extracted_range = f"{start + 1}-{actual_end + 1}"
    hint = ""
    if truncated and actual_end < total_slides - 1:
        hint = f"Content truncated. Use pages=\"{actual_end + 2}-{total_slides}\" to read the rest ({total_slides - actual_end - 1} slides remaining)."

    return ExtractResponse(
        path=str(file_path),
        type="pptx",
        content=content,
        total_pages=total_slides,
        extracted_pages=extracted_range,
        truncated=truncated,
        hint=hint,
    )


def _extract_csv(file_path: pathlib.Path, max_chars: int) -> ExtractResponse:
    """Extract data from a CSV file as a markdown table."""
    import csv

    with open(file_path, newline="", encoding="utf-8", errors="replace") as f:
        reader = csv.reader(f)
        md_rows: list[str] = []
        char_count = 0
        truncated = False
        total_rows = 0

        for ri, row in enumerate(reader):
            total_rows += 1
            cells = [c.strip() for c in row]
            line = "| " + " | ".join(cells) + " |"
            if ri == 0:
                line += "\n| " + " | ".join("---" for _ in cells) + " |"

            if char_count + len(line) > max_chars:
                truncated = True
                break

            md_rows.append(line)
            char_count += len(line)

    content = "\n".join(md_rows)
    hint = ""
    if truncated:
        hint = f"Content truncated at {max_chars} characters. File has {total_rows}+ rows."

    return ExtractResponse(
        path=str(file_path),
        type="csv",
        content=content,
        truncated=truncated,
        hint=hint,
    )


def _resolve_document_path(path_str: str) -> pathlib.Path:
    """Resolve a document path relative to workspace root, with validation."""
    p = pathlib.Path(path_str)
    if not p.is_absolute():
        if not _workspace_dir:
            raise HTTPException(status_code=400, detail="No workspace is open.")
        p = pathlib.Path(_workspace_dir) / p
    p = p.resolve()
    if not p.exists():
        raise HTTPException(status_code=404, detail=f"File not found: {path_str}")
    if not p.is_file():
        raise HTTPException(status_code=400, detail=f"Not a file: {path_str}")
    return p


# ---------------------------------------------------------------------------
# Document API endpoints
# ---------------------------------------------------------------------------

@app.get("/api/documents", response_model=list[DocumentEntry])
def list_documents() -> list[DocumentEntry]:
    """List all document and image files in the workspace.

    Walks the workspace directory and returns files whose extension matches
    a known document or image format. This is not affected by .ignore — it
    uses Python pathlib directly, so the AI can discover files that ripgrep
    would skip.
    """
    if not _workspace_dir:
        raise HTTPException(status_code=400, detail="No workspace is open.")
    root = pathlib.Path(_workspace_dir)
    if not root.exists():
        raise HTTPException(status_code=400, detail="Workspace directory not found.")

    allowed = _DOCUMENT_EXTENSIONS | _IMAGE_EXTENSIONS
    entries: list[DocumentEntry] = []

    for fp in root.rglob("*"):
        if not fp.is_file():
            continue
        # Skip hidden directories and known noise
        parts = fp.relative_to(root).parts
        if any(part.startswith(".") or part in _LS_IGNORE for part in parts[:-1]):
            continue
        ext = fp.suffix.lower()
        if ext not in allowed:
            continue
        size = fp.stat().st_size
        entries.append(DocumentEntry(
            name=fp.name,
            path=str(fp.relative_to(root)),
            type=ext.lstrip("."),
            size=size,
            size_human=_human_size(size),
        ))

    # Sort: documents first, then images; within each group alphabetical
    entries.sort(key=lambda e: (e.type in {et.lstrip(".") for et in _IMAGE_EXTENSIONS}, e.path))
    return entries


@app.post("/api/documents/extract", response_model=ExtractResponse)
def extract_document(req: ExtractRequest) -> ExtractResponse:
    """Extract content from a document file and return it as markdown text.

    Supported formats: PDF, Excel (.xlsx/.xls), Word (.docx), PowerPoint
    (.pptx), and CSV. The server enforces a hard character limit to protect
    against context window overflow.
    """
    fp = _resolve_document_path(req.path)
    ext = fp.suffix.lower()
    max_chars = min(req.max_chars, _MAX_EXTRACT_CHARS)

    if ext == ".pdf":
        return _extract_pdf(fp, max_chars, req.pages)
    elif ext in (".xlsx", ".xls"):
        return _extract_xlsx(fp, max_chars)
    elif ext in (".docx", ".doc"):
        return _extract_docx(fp, max_chars)
    elif ext in (".pptx", ".ppt"):
        return _extract_pptx(fp, max_chars, req.pages)
    elif ext == ".csv":
        return _extract_csv(fp, max_chars)
    else:
        raise HTTPException(
            status_code=400,
            detail=f"Unsupported file type: {ext}. Supported: .pdf, .xlsx, .xls, .docx, .doc, .pptx, .ppt, .csv",
        )


@app.get("/api/documents/info", response_model=DocumentInfoResponse)
def document_info(path: str = Query(..., description="Relative or absolute path to the file")) -> DocumentInfoResponse:
    """Return metadata about a document file without extracting its content.

    Useful for the AI to check page count / sheet names before deciding how
    much to extract.
    """
    fp = _resolve_document_path(path)
    ext = fp.suffix.lower()
    size = fp.stat().st_size
    resp = DocumentInfoResponse(
        path=str(fp),
        type=ext.lstrip("."),
        size=size,
        size_human=_human_size(size),
    )

    if ext == ".pdf":
        import pdfplumber
        with pdfplumber.open(fp) as pdf:
            resp.total_pages = len(pdf.pages)
    elif ext in (".xlsx", ".xls"):
        import openpyxl
        wb = openpyxl.load_workbook(fp, read_only=True, data_only=True)
        resp.total_sheets = len(wb.sheetnames)
        resp.sheet_names = wb.sheetnames
        wb.close()
    elif ext in (".pptx", ".ppt"):
        from pptx import Presentation
        prs = Presentation(fp)
        resp.total_pages = len(prs.slides)

    return resp


@app.get("/api/settings", response_model=SettingsResponse)
def get_settings() -> SettingsResponse:
    """Return current provider settings read from workspace opencode.jsonc (API key is masked)."""
    provider_id = ""
    api_key = ""
    base_url = ""
    custom_model = ""

    if _workspace_dir:
        cfg = _read_opencode_jsonc(_workspace_dir)
        providers = cfg.get("provider", {})
        if providers:
            # Take the first (and normally only) provider block written by slides-it
            provider_id = next(iter(providers))
            opts = providers[provider_id].get("options", {})
            api_key = opts.get("apiKey", "")
            base_url = opts.get("baseURL", "")
            # custom model: look for first key under "models"
            models_block = providers[provider_id].get("models", {})
            custom_model = next(iter(models_block), "")

    if len(api_key) > 12:
        masked = api_key[:8] + "..." + api_key[-4:]
    elif api_key:
        masked = "*" * len(api_key)
    else:
        masked = ""

    return SettingsResponse(
        providerID=provider_id,
        apiKeyMasked=masked,
        baseURL=base_url,
        customModel=custom_model,
    )


@app.put("/api/settings")
def save_settings(req: SettingsRequest) -> dict[str, str]:
    """
    Persist provider settings to workspace opencode.jsonc and restart OpenCode.

    - Reads the existing apiKey from opencode.jsonc if req.apiKey is empty (no change).
    - Writes the provider block (apiKey + optional baseURL) to opencode.jsonc.
    - No apiKey → removes the provider block so opencode falls back to free tier.
    - Restarts the managed opencode process so the new config takes effect immediately.
    """
    if not _workspace_dir:
        raise HTTPException(status_code=400, detail="No active workspace")

    # Resolve final key: keep existing if user left the field blank
    final_key = req.apiKey
    if not final_key:
        existing_cfg = _read_opencode_jsonc(_workspace_dir)
        providers = existing_cfg.get("provider", {})
        if req.providerID and req.providerID in providers:
            final_key = providers[req.providerID].get("options", {}).get("apiKey", "")
        elif providers:
            # If provider changed, don't carry over old key
            final_key = ""

    _write_opencode_jsonc(_workspace_dir, req.providerID, final_key, req.baseURL, req.customModel)

    # Restart opencode in a background thread so this request returns immediately.
    # The frontend polls /api/status until ready=true.
    global _opencode_proc
    if _opencode_proc is not None:
        workspace = _workspace_dir  # capture for thread closure

        def _restart() -> None:
            global _opencode_proc
            _stop_opencode()
            _opencode_proc = subprocess.Popen(
                ["opencode", "serve", "--port", str(OPENCODE_PORT)],
                cwd=workspace,
                stdout=subprocess.DEVNULL,
                stderr=subprocess.DEVNULL,
            )

        threading.Thread(target=_restart, daemon=True).start()
        return {"status": "restarting"}

    return {"status": "saved"}


@app.get("/api/file-base64")
def get_file_base64(path: str = Query(...)) -> dict[str, str]:
    """
    Read a file as raw bytes and return it base64-encoded.

    Used by the frontend to attach binary files (images, PDFs) as FileParts
    for the AI. Reading via this endpoint avoids the text-encoding corruption
    that occurs when binary data is fetched through opencode's /file/content
    endpoint (which treats all file content as UTF-8 text).

    Args:
        path: Absolute path to the file on disk.

    Returns:
        { "base64": "<base64-encoded bytes>", "mime": "<mime-type>" }
    """
    file_path = pathlib.Path(path)
    if not file_path.exists() or not file_path.is_file():
        raise HTTPException(status_code=404, detail=f"File not found: {path}")

    import base64

    mime, _ = mimetypes.guess_type(str(file_path))
    if not mime:
        mime = "application/octet-stream"

    data = file_path.read_bytes()
    b64 = base64.b64encode(data).decode("ascii")

    return {"base64": b64, "mime": mime}


@app.post("/api/upload")
async def upload_files(
    files: List[UploadFile] = File(...),
    target_dir: str = Query(default=""),
) -> dict[str, list[str]]:
    """
    Upload one or more files into the active workspace directory.

    Args:
        files: Multipart file uploads.
        target_dir: Relative sub-path within the workspace (default: workspace root).

    Returns:
        { "uploaded": ["file1.html", ...] }
    """
    if not _workspace_dir:
        raise HTTPException(status_code=400, detail="No workspace is active. Start a workspace first.")

    base = pathlib.Path(_workspace_dir)
    if target_dir:
        dest = (base / target_dir).resolve()
        # Security: must remain inside workspace
        if not str(dest).startswith(str(base)):
            raise HTTPException(status_code=400, detail="target_dir must be inside the workspace.")
    else:
        dest = base

    dest.mkdir(parents=True, exist_ok=True)

    uploaded: list[str] = []
    for upload in files:
        filename = pathlib.Path(upload.filename or "upload").name
        out_path = dest / filename
        content = await upload.read()
        out_path.write_bytes(content)
        uploaded.append(filename)

    return {"uploaded": uploaded}


@app.get("/api/file/serve")
def serve_file(path: str = Query(...)) -> Response:
    """
    Serve a workspace file for browser preview (new tab).

    Security: path must be absolute and inside the active workspace.

    Returns the raw file bytes with the correct Content-Type header so the
    browser can display images, PDFs, HTML, etc. directly.
    """
    if not _workspace_dir:
        raise HTTPException(status_code=400, detail="No active workspace")

    file_path = pathlib.Path(path).resolve()
    workspace = pathlib.Path(_workspace_dir).resolve()

    # Security: reject any path outside the workspace
    try:
        file_path.relative_to(workspace)
    except ValueError:
        raise HTTPException(status_code=403, detail="Path is outside the workspace")

    if not file_path.exists() or not file_path.is_file():
        raise HTTPException(status_code=404, detail=f"File not found: {path}")

    mime, _ = mimetypes.guess_type(str(file_path))
    if not mime:
        mime = "application/octet-stream"

    data = file_path.read_bytes()
    return Response(content=data, media_type=mime)


@app.put("/api/file/content")
def update_file_content(req: FileContentRequest) -> dict[str, str]:
    """
    Write text content to a workspace file.

    Used by the PreviewPanel Save button to persist inline edits made in
    the browser back to the original HTML file on disk.

    Security: path must be absolute and inside the active workspace.
    """
    if not _workspace_dir:
        raise HTTPException(status_code=400, detail="No active workspace")

    file_path = pathlib.Path(req.path).resolve()
    workspace = pathlib.Path(_workspace_dir).resolve()

    try:
        file_path.relative_to(workspace)
    except ValueError:
        raise HTTPException(status_code=403, detail="Path is outside the workspace")

    if not file_path.exists() or not file_path.is_file():
        raise HTTPException(status_code=404, detail=f"File not found: {req.path}")

    file_path.write_text(req.content, encoding="utf-8")
    return {"path": str(file_path), "status": "saved"}


@app.post("/api/export/bundle")
def export_bundle(req: BundleRequest) -> dict[str, str]:
    """
    Bundle an HTML file with all local images inlined as base64 data URIs.

    Scans the HTML for ``<img src="...">`` and CSS ``url(...)`` references
    to local files, reads each file, and replaces the reference with an
    inline ``data:<mime>;base64,...`` URI.  The result is a fully
    self-contained HTML string that can be shared without losing images.

    External URLs (http://, https://, data:) are left unchanged.

    Returns:
        { "content": "<bundled HTML>", "filename": "<name>-bundled.html" }
    """
    if not _workspace_dir:
        raise HTTPException(status_code=400, detail="No active workspace")

    file_path = pathlib.Path(req.path).resolve()
    workspace = pathlib.Path(_workspace_dir).resolve()

    try:
        file_path.relative_to(workspace)
    except ValueError:
        raise HTTPException(status_code=403, detail="Path is outside the workspace")

    if not file_path.exists() or not file_path.is_file():
        raise HTTPException(status_code=404, detail=f"File not found: {req.path}")

    html = file_path.read_text(encoding="utf-8")
    html_dir = file_path.parent

    def _inline_local_ref(match: re.Match[str]) -> str:
        """Replace a local file reference with its base64 data URI."""
        prefix = match.group(1)   # e.g. 'src="' or "url('"
        ref = match.group(2)      # the path/URL
        suffix = match.group(3)   # e.g. '"' or "')"

        # Skip external URLs and data URIs
        if ref.startswith(("http://", "https://", "data:", "//", "#")):
            return match.group(0)

        # Resolve relative to the HTML file's directory
        asset_path = (html_dir / ref).resolve()

        # Security: must be inside workspace
        try:
            asset_path.relative_to(workspace)
        except ValueError:
            return match.group(0)

        if not asset_path.exists() or not asset_path.is_file():
            return match.group(0)

        mime, _ = mimetypes.guess_type(str(asset_path))
        if not mime:
            return match.group(0)

        raw = asset_path.read_bytes()
        b64 = base64.b64encode(raw).decode("ascii")
        return f"{prefix}data:{mime};base64,{b64}{suffix}"

    # Match <img src="...">, <source src="...">, etc.
    html = re.sub(
        r'''(src\s*=\s*["'])([^"']+)(["'])''',
        _inline_local_ref,
        html,
    )

    # Match CSS url("...") / url('...') / url(...)
    html = re.sub(
        r'''(url\(\s*["']?)([^"')]+)(["']?\s*\))''',
        _inline_local_ref,
        html,
    )

    stem = file_path.stem
    filename = f"{stem}-bundled.html"

    return {"content": html, "filename": filename}


@app.put("/api/file/rename")
def rename_file(req: FileRenameRequest) -> dict[str, str]:
    """
    Rename a file or directory inside the active workspace.

    Args:
        path:     Absolute path of the existing file/directory.
        new_name: New filename only (no directory components).

    Returns:
        { "path": "<new absolute path>" }
    """
    if not _workspace_dir:
        raise HTTPException(status_code=400, detail="No active workspace")

    src = pathlib.Path(req.path).resolve()
    workspace = pathlib.Path(_workspace_dir).resolve()

    try:
        src.relative_to(workspace)
    except ValueError:
        raise HTTPException(status_code=403, detail="Path is outside the workspace")

    if not src.exists():
        raise HTTPException(status_code=404, detail=f"Not found: {req.path}")

    # Reject new_name with path separators to prevent directory traversal
    new_name = req.new_name.strip()
    if not new_name or "/" in new_name or "\\" in new_name or new_name in (".", ".."):
        raise HTTPException(status_code=400, detail="Invalid new_name")

    dst = src.parent / new_name
    if dst.exists():
        raise HTTPException(status_code=409, detail=f"Already exists: {new_name}")

    src.rename(dst)
    return {"path": str(dst)}


@app.delete("/api/file")
def delete_file(path: str = Query(...)) -> dict[str, str]:
    """
    Delete a file or directory (recursively) inside the active workspace.

    Security: path must be absolute and inside the workspace.
    """
    if not _workspace_dir:
        raise HTTPException(status_code=400, detail="No active workspace")

    target = pathlib.Path(path).resolve()
    workspace = pathlib.Path(_workspace_dir).resolve()

    try:
        target.relative_to(workspace)
    except ValueError:
        raise HTTPException(status_code=403, detail="Path is outside the workspace")

    if not target.exists():
        raise HTTPException(status_code=404, detail=f"Not found: {path}")

    if target.is_dir():
        shutil.rmtree(target)
    else:
        target.unlink()

    return {"path": str(target), "status": "deleted"}


@app.post("/api/file")
def create_file(req: FileCreateRequest) -> dict[str, str]:
    """
    Create an empty file inside the workspace root.

    Args:
        name: Filename only (no path separators).

    Returns:
        { "path": "<absolute path of new file>" }
    """
    if not _workspace_dir:
        raise HTTPException(status_code=400, detail="No active workspace")

    name = req.name.strip()
    if not name or "/" in name or "\\" in name or name in (".", ".."):
        raise HTTPException(status_code=400, detail="Invalid file name")

    target = pathlib.Path(_workspace_dir).resolve() / name
    if target.exists():
        raise HTTPException(status_code=409, detail=f"Already exists: {name}")

    target.touch()
    return {"path": str(target), "status": "created"}


@app.post("/api/mkdir")
def create_directory(req: DirCreateRequest) -> dict[str, str]:
    """
    Create a directory inside the workspace root.

    Args:
        name: Directory name only (no path separators).

    Returns:
        { "path": "<absolute path of new directory>" }
    """
    if not _workspace_dir:
        raise HTTPException(status_code=400, detail="No active workspace")

    name = req.name.strip()
    if not name or "/" in name or "\\" in name or name in (".", ".."):
        raise HTTPException(status_code=400, detail="Invalid directory name")

    target = pathlib.Path(_workspace_dir).resolve() / name
    if target.exists():
        raise HTTPException(status_code=409, detail=f"Already exists: {name}")

    target.mkdir(parents=True)
    return {"path": str(target), "status": "created"}


# ---------------------------------------------------------------------------
# Replay — infinite context
# ---------------------------------------------------------------------------

_replay_log = logging.getLogger("slides-it.replay")

# Patterns that signal the LLM provider rejected the request because the
# conversation exceeded the model's context window.
_OVERFLOW_PATTERNS: list[re.Pattern[str]] = [
    re.compile(r"context.{0,20}(too long|too large|exceed|overflow)", re.I),
    re.compile(r"(max|maximum).{0,20}(token|context).{0,20}(limit|exceed)", re.I),
    re.compile(r"(token|context).{0,20}(limit|budget).{0,20}(exceed|reach)", re.I),
    re.compile(r"request too large", re.I),
    re.compile(r"prompt is too long", re.I),
    re.compile(r"content_too_large", re.I),
    re.compile(r"input.*too long", re.I),
    re.compile(r"max_tokens_exceeded", re.I),
    re.compile(r"context_length_exceeded", re.I),
    re.compile(r"token limit", re.I),
]


def _is_context_overflow(error_msg: str) -> bool:
    """Return True if *error_msg* looks like a context-window overflow."""
    return any(pat.search(error_msg) for pat in _OVERFLOW_PATTERNS)


def _opencode_get(path: str) -> list | dict:
    """GET request to OpenCode server, return parsed JSON."""
    url = f"http://localhost:{OPENCODE_PORT}{path}"
    req = urllib.request.Request(url)
    req.add_header("Content-Type", "application/json")
    with urllib.request.urlopen(req, timeout=30) as resp:
        return json.loads(resp.read())


def _opencode_post(
    path: str, body: dict, *, expect_json: bool = True
) -> dict | list | bool | str | None:
    """POST request to OpenCode server."""
    url = f"http://localhost:{OPENCODE_PORT}{path}"
    data = json.dumps(body).encode()
    req = urllib.request.Request(url, data=data, method="POST")
    req.add_header("Content-Type", "application/json")
    with urllib.request.urlopen(req, timeout=60) as resp:
        raw = resp.read()
        if not raw or resp.status == 204:
            return None
        if expect_json:
            return json.loads(raw)
        return raw.decode(errors="replace")


def _replay_summarize(session_id: str, provider_id: str, model_id: str) -> bool:
    """Call POST /session/{id}/summarize and return success."""
    body = {"providerID": provider_id, "modelID": model_id}
    resp = _opencode_post(f"/session/{session_id}/summarize", body)
    return resp is True or resp == "true"


def _replay_extract_summary(session_id: str) -> str:
    """
    Read back messages and find the summary generated by ``summarize``.

    The summarize call appends a user message with a ``compaction`` part
    followed by an assistant message whose ``text`` parts contain the summary.
    """
    messages: list[dict] = _opencode_get(f"/session/{session_id}/message")

    # Walk from the end — find the assistant message right after compaction.
    for i in range(len(messages) - 1, -1, -1):
        msg = messages[i]
        if msg["info"]["role"] != "assistant":
            continue
        # Check if the previous message is a compaction marker
        if i > 0:
            prev_parts = messages[i - 1].get("parts", [])
            if any(p.get("type") == "compaction" for p in prev_parts):
                texts = [
                    p["text"]
                    for p in msg.get("parts", [])
                    if p.get("type") == "text" and p.get("text")
                ]
                if texts:
                    return "\n\n".join(texts)

    # Fallback: grab the last assistant text
    for i in range(len(messages) - 1, -1, -1):
        msg = messages[i]
        if msg["info"]["role"] == "assistant":
            texts = [
                p["text"]
                for p in msg.get("parts", [])
                if p.get("type") == "text" and p.get("text")
            ]
            if texts:
                return "\n\n".join(texts)

    return ""


def _replay_create_child(parent_id: str) -> dict:
    """Create a new session with parentID pointing to the old one."""
    return _opencode_post("/session", {"parentID": parent_id, "title": "slides-it"})


def _replay_inject_context(
    session_id: str, summary: str, system_prompt: str | None = None
) -> None:
    """Inject compacted context into the new session without triggering AI reply."""
    context = (
        "[Context from previous conversation — compacted]\n\n"
        f"{summary}\n\n"
        "[End of compacted context — continue the conversation from here]"
    )
    body: dict = {
        "noReply": True,
        "parts": [{"type": "text", "text": context}],
    }
    if system_prompt:
        body["system"] = system_prompt
    _opencode_post(f"/session/{session_id}/prompt_async", body, expect_json=False)


def _do_replay(
    session_id: str,
    provider_id: str,
    model_id: str,
    system_prompt: str | None = None,
) -> dict[str, str]:
    """
    Execute the full replay flow:

    1. Summarize the current session (compaction)
    2. Extract the summary text from the compaction messages
    3. Create a new child session
    4. Inject the summary into the child session as context

    Returns dict with ``new_session_id``, ``parent_session_id``, ``summary``.
    """
    _replay_log.info("replay: starting for session %s", session_id)

    # 1. Summarize
    _replay_log.info("replay: step 1 — summarize")
    ok = _replay_summarize(session_id, provider_id, model_id)
    if not ok:
        raise RuntimeError(f"summarize returned failure for session {session_id}")

    # 2. Extract summary
    _replay_log.info("replay: step 2 — extract summary")
    summary = _replay_extract_summary(session_id)
    if not summary:
        raise RuntimeError(f"could not extract summary from session {session_id}")
    _replay_log.info("replay: summary %d chars", len(summary))

    # 3. Create child session
    _replay_log.info("replay: step 3 — create child session")
    new_session = _replay_create_child(session_id)
    new_id = new_session["id"]
    _replay_log.info("replay: new session %s (parent=%s)", new_id, session_id)

    # 4. Inject context
    _replay_log.info("replay: step 4 — inject context into %s", new_id)
    _replay_inject_context(new_id, summary, system_prompt)

    _replay_log.info("replay: complete — %s → %s", session_id, new_id)
    return {
        "new_session_id": new_id,
        "parent_session_id": session_id,
        "summary": summary,
    }


@app.post("/api/replay")
def replay(req: ReplayRequest) -> dict[str, str]:
    """
    Compact the current session and continue in a new child session.

    The frontend calls this when context overflows (automatically) or
    when the user clicks a "Compact" button (manually).

    If ``provider_id`` / ``model_id`` are empty, the active model from
    the slides-it config is used.
    """
    if not _workspace_dir:
        raise HTTPException(status_code=400, detail="No active workspace")

    # Resolve provider/model — fall back to the active model in config
    provider_id = req.provider_id
    model_id = req.model_id
    if not provider_id or not model_id:
        dm = DesignManager()
        current = dm.get_model()  # e.g. "anthropic/claude-sonnet-4-6"
        if current and "/" in current:
            provider_id = provider_id or current.split("/", 1)[0]
            model_id = model_id or current.split("/", 1)[1]
    if not provider_id or not model_id:
        raise HTTPException(
            status_code=400,
            detail="Cannot determine model for summarization. Set provider_id and model_id, or configure an active model.",
        )

    # Optionally build the system prompt so it is injected into the new session
    system_prompt: str | None = None
    try:
        dm = DesignManager()
        active_design = dm.get_active()
        im = IndustryManager()
        active_industry = im.get_active()
        system_prompt = dm.build_prompt(active_design, active_industry)
    except Exception:
        pass  # non-fatal — the new session will still work without a system prompt

    try:
        result = _do_replay(req.session_id, provider_id, model_id, system_prompt)
    except RuntimeError as e:
        raise HTTPException(status_code=502, detail=str(e)) from e
    except urllib.error.HTTPError as e:
        body = e.read().decode(errors="replace")[:500]
        raise HTTPException(
            status_code=502, detail=f"OpenCode error: HTTP {e.code} — {body}"
        ) from e
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Replay failed: {e}") from e

    return result


@app.post("/api/replay/check")
def replay_check(body: dict[str, str]) -> dict[str, bool]:
    """
    Check if an error message indicates a context overflow.

    Body: ``{ "error": "..." }``
    Returns: ``{ "is_overflow": true/false }``

    The frontend calls this from the session.error SSE handler to decide
    whether to trigger an automatic replay.
    """
    error_msg = body.get("error", "")
    return {"is_overflow": _is_context_overflow(error_msg)}


# ---------------------------------------------------------------------------
# Static frontend (served last so API routes take precedence)
# ---------------------------------------------------------------------------

def mount_frontend(dist_path: pathlib.Path) -> None:
    """Mount the built frontend at /. Call after app is created."""
    if dist_path.exists():
        app.mount("/", StaticFiles(directory=str(dist_path), html=True), name="frontend")


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

_OPENCODE_AGENTS_MD = pathlib.Path.home() / ".config" / "opencode" / "AGENTS.md"
_SLIDES_IT_MARKER = "<!-- slides-it managed -->"


def _cleanup_stale_agents_md() -> None:
    """
    Remove any slides-it block left in ~/.config/opencode/AGENTS.md by an older
    version that used write_rules(). slides-it no longer writes to that file;
    system prompts are passed per-message via the `system` API field instead.
    Safe to call on every startup — no-ops if the file doesn't exist or has no
    slides-it marker.
    """
    if not _OPENCODE_AGENTS_MD.exists():
        return
    text = _OPENCODE_AGENTS_MD.read_text(encoding="utf-8")
    if _SLIDES_IT_MARKER not in text:
        return
    parts = text.split(_SLIDES_IT_MARKER)
    # Preserve any user content that existed outside the slides-it block
    tail = parts[-1].strip() if len(parts) > 2 else ""
    if tail:
        _OPENCODE_AGENTS_MD.write_text(tail, encoding="utf-8")
    else:
        _OPENCODE_AGENTS_MD.unlink(missing_ok=True)


def _is_opencode_healthy() -> bool:
    """Return True if opencode is already responding on OPENCODE_PORT."""
    try:
        with urllib.request.urlopen(
            f"http://localhost:{OPENCODE_PORT}/global/health", timeout=1
        ) as resp:
            data = json.loads(resp.read())
            return bool(data.get("healthy", False))
    except Exception:
        return False


def _read_opencode_jsonc(workspace: str) -> dict:
    """
    Read and parse the slides-it provider config from the workspace.

    Reads opencode.json first (the file opencode actually loads as project config).
    Falls back to opencode.jsonc for users who had an older slides-it version.

    Returns an empty dict if no file exists or cannot be parsed.
    """
    import re
    # opencode loads opencode.json as the project config — use that as primary
    cfg_path = pathlib.Path(workspace) / "opencode.json"
    if not cfg_path.exists():
        cfg_path = pathlib.Path(workspace) / "opencode.jsonc"
    if not cfg_path.exists():
        return {}
    try:
        text = cfg_path.read_text(encoding="utf-8")
        try:
            return json.loads(text)
        except json.JSONDecodeError:
            # Strip only lines whose first non-whitespace chars are "//"
            stripped = re.sub(r"(?m)^\s*//.*$", "", text)
            return json.loads(stripped)
    except Exception:
        return {}


def _write_opencode_jsonc(
    workspace: str,
    provider_id: str,
    api_key: str,
    base_url: str,
    custom_model: str,
) -> None:
    """
    Write provider config into <workspace>/opencode.json.

    Uses opencode.json (not .jsonc) because that is the filename opencode
    recognises as the project-level config. Writing to .jsonc would be silently
    ignored by opencode when a .json file is also present.
    """
    import re
    # Write to opencode.json — the file opencode actually loads as project config
    cfg_path = pathlib.Path(workspace) / "opencode.json"

    # Parse existing file (comment-stripping for safety)
    cfg: dict = {}
    if cfg_path.exists():
        try:
            text = cfg_path.read_text(encoding="utf-8")
            try:
                cfg = json.loads(text)
            except json.JSONDecodeError:
                stripped = re.sub(r"(?m)^\s*//.*$", "", text)
                cfg = json.loads(stripped)
        except Exception:
            cfg = {}

    cfg.setdefault("$schema", "https://opencode.ai/config.json")

    # ── MCP: always enable open-webSearch for AI web search capability ──
    cfg["mcp"] = {"web-search": _MCP_WEB_SEARCH_BLOCK}

    provider_section: dict = cfg.get("provider", {})

    # Remove stale provider blocks from previous provider (clean slate on switch)
    if provider_id and provider_section:
        stale = [k for k in list(provider_section.keys()) if k != provider_id]
        for k in stale:
            del provider_section[k]

    if api_key and provider_id:
        options: dict = {"apiKey": api_key, "setCacheKey": True}
        if base_url:
            options["baseURL"] = base_url

        if provider_id in _NATIVE_PROVIDERS:
            # Native provider: just options block
            block: dict = {"options": options}
        else:
            # Custom OpenAI-compatible provider: needs npm + name + options
            block = {
                "npm": "@ai-sdk/openai-compatible",
                "name": provider_id,
                "options": options,
            }
            if custom_model:
                block["models"] = {custom_model: {"name": custom_model}}

        provider_section[provider_id] = block
        cfg["provider"] = provider_section
    else:
        # No key → remove the provider block entirely (fall back to free tier)
        if provider_id in provider_section:
            del provider_section[provider_id]
        if provider_section:
            cfg["provider"] = provider_section
        elif "provider" in cfg:
            del cfg["provider"]

    cfg_path.write_text(json.dumps(cfg, indent=2), encoding="utf-8")


def _ensure_mcp_config(workspace: str) -> None:
    """
    Ensure the workspace opencode.json contains the MCP web-search config
    and that all provider options include setCacheKey: true.

    Called on every workspace start so that web search and prompt caching are
    available even if the user has never opened Settings.  Reads the existing
    file, injects missing config, and writes back only when something changed.
    """
    import re

    cfg_path = pathlib.Path(workspace) / "opencode.json"
    cfg: dict = {}
    if cfg_path.exists():
        try:
            text = cfg_path.read_text(encoding="utf-8")
            try:
                cfg = json.loads(text)
            except json.JSONDecodeError:
                stripped = re.sub(r"(?m)^\s*//.*$", "", text)
                cfg = json.loads(stripped)
        except Exception:
            cfg = {}

    changed = False

    # ── MCP: ensure web-search is present ──
    mcp = cfg.get("mcp", {})
    if "web-search" not in mcp:
        mcp["web-search"] = _MCP_WEB_SEARCH_BLOCK
        cfg["mcp"] = mcp
        changed = True

    # ── Provider: ensure setCacheKey is present in all provider options ──
    providers = cfg.get("provider", {})
    for _pid, pblock in providers.items():
        if not isinstance(pblock, dict):
            continue
        opts = pblock.get("options", {})
        if not isinstance(opts, dict):
            continue
        if opts.get("setCacheKey") is not True:
            opts["setCacheKey"] = True
            pblock["options"] = opts
            changed = True

    if changed:
        cfg_path.write_text(json.dumps(cfg, indent=2), encoding="utf-8")


def _stop_opencode() -> None:
    global _opencode_proc
    if _opencode_proc and _opencode_proc.poll() is None:
        _opencode_proc.terminate()
        try:
            _opencode_proc.wait(timeout=5)
        except subprocess.TimeoutExpired:
            _opencode_proc.kill()
    _opencode_proc = None
    # Also kill any stray process occupying the opencode port
    _free_port(OPENCODE_PORT)


def _free_port(port: int) -> None:
    """Kill any process listening on the given port (macOS / Linux)."""
    import sys
    if sys.platform == "win32":
        return
    try:
        result = subprocess.run(
            ["lsof", "-ti", f":{port}"],
            capture_output=True,
            text=True,
        )
        for pid in result.stdout.strip().splitlines():
            pid = pid.strip()
            if pid:
                subprocess.run(["kill", "-9", pid], check=False)
    except FileNotFoundError:
        pass


# ---------------------------------------------------------------------------
# Entrypoint (for direct invocation / dev)
# ---------------------------------------------------------------------------

def run(port: int = 3000, frontend_dist: pathlib.Path | None = None) -> None:
    if frontend_dist:
        mount_frontend(frontend_dist)
    uvicorn.run(app, host="127.0.0.1", port=port, log_level="warning")
